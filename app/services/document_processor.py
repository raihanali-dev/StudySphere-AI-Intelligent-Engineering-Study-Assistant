"""Document processing service."""
import os
from typing import List, Tuple
from pathlib import Path


class DocumentProcessor:
    """Handle document processing and text extraction."""

    CHUNK_SIZE = 500
    CHUNK_OVERLAP = 100

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        import pdfplumber

        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        from docx import Document

        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from DOCX: {str(e)}")

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        from pptx import Presentation

        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from PPTX: {str(e)}")

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Error extracting text from TXT: {str(e)}")

    @classmethod
    def extract_text(cls, file_path: str, file_type: str) -> str:
        """Extract text from any supported file type."""
        if file_type.lower() == "pdf":
            return cls.extract_text_from_pdf(file_path)
        elif file_type.lower() == "docx":
            return cls.extract_text_from_docx(file_path)
        elif file_type.lower() == "pptx":
            return cls.extract_text_from_pptx(file_path)
        elif file_type.lower() == "txt":
            return cls.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    @classmethod
    def chunk_text(cls, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """Split text into overlapping chunks."""
        chunk_size = chunk_size or cls.CHUNK_SIZE
        overlap = overlap or cls.CHUNK_OVERLAP

        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunks.append(text[i : i + chunk_size])

        return chunks

    @classmethod
    def clean_text(cls, text: str) -> str:
        """Clean and normalize text."""
        # Remove extra whitespace
        text = " ".join(text.split())
        # Remove special characters (keep essential ones)
        return text
